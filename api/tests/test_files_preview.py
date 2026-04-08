"""Tests for the /files/preview and /files/raw endpoints."""

from __future__ import annotations

import tempfile
from pathlib import Path
from unittest.mock import patch

import pytest


# --- /files/preview rich handlers ---


@pytest.mark.asyncio
async def test_preview_rejects_path_traversal(client):
    resp = await client.get("/api/files/preview?path=../../etc/passwd")
    assert resp.status_code == 403


@pytest.mark.asyncio
async def test_preview_rejects_missing(client):
    resp = await client.get("/api/files/preview?path=no-such-file.docx")
    assert resp.status_code == 404


@pytest.mark.asyncio
async def test_preview_rejects_directory(client):
    with tempfile.TemporaryDirectory() as tmpdir:
        tmppath = Path(tmpdir)
        (tmppath / "subdir").mkdir()
        with patch("routers.projects.TORIOS_DIR", tmppath):
            resp = await client.get("/api/files/preview?path=subdir")
    assert resp.status_code == 400


@pytest.mark.asyncio
async def test_preview_unknown_extension_rejected(client):
    """Unknown extensions should return 400 from the preview endpoint.

    The client should fall back to /files/read or /files/raw.
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        tmppath = Path(tmpdir)
        (tmppath / "notes.md").write_text("# hello")
        with patch("routers.projects.TORIOS_DIR", tmppath):
            resp = await client.get("/api/files/preview?path=notes.md")
    assert resp.status_code == 400
    assert "no rich preview" in resp.json()["detail"].lower()


@pytest.mark.asyncio
async def test_preview_csv(client):
    with tempfile.TemporaryDirectory() as tmpdir:
        tmppath = Path(tmpdir)
        csv_path = tmppath / "data.csv"
        csv_path.write_text("name,age\nAlice,30\nBob,25\n")

        with patch("routers.projects.TORIOS_DIR", tmppath):
            resp = await client.get("/api/files/preview?path=data.csv")

    assert resp.status_code == 200
    data = resp.json()
    assert data["kind"] == "spreadsheet"
    assert len(data["sheets"]) == 1
    sheet = data["sheets"][0]
    assert sheet["rows"][0] == ["name", "age"]
    assert sheet["rows"][1] == ["Alice", "30"]
    assert sheet["total_rows"] == 3
    assert sheet["truncated"] is False


@pytest.mark.asyncio
async def test_preview_csv_truncates_large_file(client):
    """A CSV larger than the row cap should be truncated and flagged."""
    from routers.files import MAX_PREVIEW_ROWS

    with tempfile.TemporaryDirectory() as tmpdir:
        tmppath = Path(tmpdir)
        csv_path = tmppath / "big.csv"
        lines = ["col\n"] + [f"row{i}\n" for i in range(MAX_PREVIEW_ROWS + 50)]
        csv_path.write_text("".join(lines))

        with patch("routers.projects.TORIOS_DIR", tmppath):
            resp = await client.get("/api/files/preview?path=big.csv")

    assert resp.status_code == 200
    sheet = resp.json()["sheets"][0]
    assert len(sheet["rows"]) == MAX_PREVIEW_ROWS
    assert sheet["truncated"] is True
    assert sheet["total_rows"] > MAX_PREVIEW_ROWS


@pytest.mark.asyncio
async def test_preview_xlsx(client):
    from openpyxl import Workbook

    with tempfile.TemporaryDirectory() as tmpdir:
        tmppath = Path(tmpdir)
        xlsx_path = tmppath / "workbook.xlsx"

        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws.append(["name", "score"])
        ws.append(["Alice", 95])
        ws.append(["Bob", 82])
        wb.create_sheet("Extras").append(["note"])
        wb.save(xlsx_path)
        wb.close()

        with patch("routers.projects.TORIOS_DIR", tmppath):
            resp = await client.get("/api/files/preview?path=workbook.xlsx")

    assert resp.status_code == 200
    data = resp.json()
    assert data["kind"] == "spreadsheet"
    names = [s["name"] for s in data["sheets"]]
    assert "Sheet1" in names
    assert "Extras" in names
    sheet1 = next(s for s in data["sheets"] if s["name"] == "Sheet1")
    assert sheet1["rows"][0] == ["name", "score"]
    assert sheet1["rows"][1] == ["Alice", 95]


@pytest.mark.asyncio
async def test_preview_docx(client):
    from docx import Document

    with tempfile.TemporaryDirectory() as tmpdir:
        tmppath = Path(tmpdir)
        docx_path = tmppath / "letter.docx"

        doc = Document()
        doc.add_heading("Hello World", level=1)
        doc.add_paragraph("This is a test paragraph.")
        doc.save(docx_path)

        with patch("routers.projects.TORIOS_DIR", tmppath):
            resp = await client.get("/api/files/preview?path=letter.docx")

    assert resp.status_code == 200
    data = resp.json()
    assert data["kind"] == "docx"
    assert "Hello World" in data["html"]
    assert "test paragraph" in data["html"]


@pytest.mark.asyncio
async def test_preview_pptx(client):
    from pptx import Presentation

    with tempfile.TemporaryDirectory() as tmpdir:
        tmppath = Path(tmpdir)
        pptx_path = tmppath / "deck.pptx"

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        if title is not None:
            title.text = "Welcome"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "Intro body text"
        prs.save(pptx_path)

        with patch("routers.projects.TORIOS_DIR", tmppath):
            resp = await client.get("/api/files/preview?path=deck.pptx")

    assert resp.status_code == 200
    data = resp.json()
    assert data["kind"] == "slides"
    assert len(data["slides"]) == 1
    slide = data["slides"][0]
    assert slide["index"] == 1
    combined = (slide["title"] or "") + (slide["body"] or "")
    assert "Welcome" in combined


@pytest.mark.asyncio
async def test_preview_pdf(client):
    """Create a minimal PDF on disk and verify pypdf can open it."""
    with tempfile.TemporaryDirectory() as tmpdir:
        tmppath = Path(tmpdir)
        pdf_path = tmppath / "doc.pdf"

        # Minimal valid PDF, one blank page.
        pdf_bytes = (
            b"%PDF-1.4\n"
            b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
            b"2 0 obj<</Type/Pages/Count 1/Kids[3 0 R]>>endobj\n"
            b"3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]>>endobj\n"
            b"xref\n0 4\n0000000000 65535 f \n"
            b"0000000010 00000 n \n0000000053 00000 n \n0000000098 00000 n \n"
            b"trailer<</Size 4/Root 1 0 R>>\nstartxref\n149\n%%EOF\n"
        )
        pdf_path.write_bytes(pdf_bytes)

        with patch("routers.projects.TORIOS_DIR", tmppath):
            resp = await client.get("/api/files/preview?path=doc.pdf")

    assert resp.status_code == 200
    data = resp.json()
    assert data["kind"] == "pdf"
    assert data["total_pages"] == 1
    assert len(data["pages"]) == 1


# --- /files/raw endpoint ---


@pytest.mark.asyncio
async def test_raw_png(client):
    import base64 as b64

    with tempfile.TemporaryDirectory() as tmpdir:
        tmppath = Path(tmpdir)
        png = b64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAAC0lEQVQI12NgAAIABQAB"
            "Nl7BcQAAAABJRU5ErkJggg=="
        )
        (tmppath / "pixel.png").write_bytes(png)

        with patch("routers.projects.TORIOS_DIR", tmppath):
            resp = await client.get("/api/files/raw?path=pixel.png")

    assert resp.status_code == 200
    assert resp.headers["content-type"].startswith("image/png")
    assert resp.content == png


@pytest.mark.asyncio
async def test_raw_rejects_path_traversal(client):
    resp = await client.get("/api/files/raw?path=../../etc/passwd")
    assert resp.status_code == 403


@pytest.mark.asyncio
async def test_raw_rejects_directory(client):
    with tempfile.TemporaryDirectory() as tmpdir:
        tmppath = Path(tmpdir)
        (tmppath / "folder").mkdir()
        with patch("routers.projects.TORIOS_DIR", tmppath):
            resp = await client.get("/api/files/raw?path=folder")
    assert resp.status_code == 400
