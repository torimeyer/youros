"""Tests for the /files/beautify-deck endpoint.

These tests mock the Anthropic client so we never hit the real API.
"""

from __future__ import annotations

import tempfile
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import AsyncMock, patch

import pytest


def _make_pptx(path: Path, slide_count: int = 2) -> None:
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[0]
    for i in range(slide_count):
        slide = prs.slides.add_slide(layout)
        title = slide.shapes.title
        if title is not None:
            title.text = f"Slide {i + 1} title"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = f"Body text for slide {i + 1}"
    prs.save(path)


def _fake_model_response(
    layout: str = "title_and_content",
    title: str = "Welcome home",
    body: list | dict | None = None,
    accent: str = "#22c55e",
) -> SimpleNamespace:
    """Build a shape that mimics the anthropic SDK response."""
    if body is None:
        body = ["Clear first point", "Clear second point"]
    payload = {
        "layout": layout,
        "title": title,
        "body": body,
        "accent_color": accent,
        "font_pairing": {"heading": "Inter", "body": "Source Sans"},
        "notes": "Pause for effect after the opening line.",
        "rationale": "Tighter wording and a friendlier accent color.",
    }
    import json as _json

    text_block = SimpleNamespace(type="text", text=_json.dumps(payload))
    return SimpleNamespace(content=[text_block])


class _FakeMessages:
    def __init__(self):
        self.create = AsyncMock(return_value=_fake_model_response())


class _FakeClient:
    def __init__(self, *args, **kwargs):
        self.messages = _FakeMessages()


@pytest.mark.asyncio
async def test_beautify_happy_path(client):
    """Parse a real pptx, mock Claude, confirm the response shape."""
    with tempfile.TemporaryDirectory() as tmpdir:
        tmppath = Path(tmpdir)
        pptx_path = tmppath / "hello.pptx"
        _make_pptx(pptx_path, slide_count=2)

        with patch("routers.projects._projects_root", return_value=tmppath), patch(
            "services.ai_backend.get_ai_client", AsyncMock(return_value=_FakeClient())
        ), patch.dict(
            "routers.beautify._cache", {}, clear=True
        ):
            resp = await client.post(
                "/api/files/beautify-deck", json={"path": "hello.pptx"}
            )

    assert resp.status_code == 200, resp.text
    data = resp.json()
    assert data["name"] == "hello.pptx"
    assert data["slide_count"] == 2
    assert len(data["originals"]) == 2
    assert len(data["beautified"]) == 2

    first = data["beautified"][0]
    assert first["layout"] == "title_and_content"
    assert first["title"] == "Welcome home"
    assert first["body"] == ["Clear first point", "Clear second point"]
    assert first["accent_color"] == "#22c55e"
    assert first["font_pairing"]["heading"] == "Inter"
    assert first["font_pairing"]["body"] == "Source Sans"
    assert first["rationale"]
    assert first["index"] == 1


@pytest.mark.asyncio
async def test_beautify_rejects_non_pptx(client):
    with tempfile.TemporaryDirectory() as tmpdir:
        tmppath = Path(tmpdir)
        (tmppath / "notes.md").write_text("# hi")
        with patch("routers.projects._projects_root", return_value=tmppath):
            resp = await client.post(
                "/api/files/beautify-deck", json={"path": "notes.md"}
            )
    assert resp.status_code == 400
    assert "powerpoint" in resp.json()["detail"].lower() or ".pptx" in resp.json()["detail"]


@pytest.mark.asyncio
async def test_beautify_rejects_oversized_deck(client):
    """A deck over the cap should fail with a friendly message."""
    from routers.beautify import MAX_BEAUTIFY_SLIDES

    with tempfile.TemporaryDirectory() as tmpdir:
        tmppath = Path(tmpdir)
        pptx_path = tmppath / "huge.pptx"
        _make_pptx(pptx_path, slide_count=MAX_BEAUTIFY_SLIDES + 1)

        with patch("routers.projects._projects_root", return_value=tmppath):
            resp = await client.post(
                "/api/files/beautify-deck", json={"path": "huge.pptx"}
            )

    assert resp.status_code == 400
    detail = resp.json()["detail"].lower()
    assert "polish" in detail or "shorter" in detail
    assert str(MAX_BEAUTIFY_SLIDES) in resp.json()["detail"]


@pytest.mark.asyncio
async def test_beautify_rejects_path_traversal(client):
    resp = await client.post(
        "/api/files/beautify-deck", json={"path": "../../etc/passwd"}
    )
    assert resp.status_code == 403


@pytest.mark.asyncio
async def test_beautify_cache_uses_mtime(client):
    """Modifying the file should bust the cache and re-run the model."""
    with tempfile.TemporaryDirectory() as tmpdir:
        tmppath = Path(tmpdir)
        pptx_path = tmppath / "deck.pptx"
        _make_pptx(pptx_path, slide_count=1)

        fake_client = _FakeClient()
        # The same fake instance is used on both calls because we patch the
        # class to always return it.
        with patch("routers.projects._projects_root", return_value=tmppath), patch(
            "services.ai_backend.get_ai_client", AsyncMock(return_value=fake_client)
        ), patch.dict(
            "routers.beautify._cache", {}, clear=True
        ):
            resp1 = await client.post(
                "/api/files/beautify-deck", json={"path": "deck.pptx"}
            )
            assert resp1.status_code == 200
            assert resp1.json()["cached"] is False
            first_calls = fake_client.messages.create.call_count

            # Second call on the unchanged file uses the cache.
            resp2 = await client.post(
                "/api/files/beautify-deck", json={"path": "deck.pptx"}
            )
            assert resp2.status_code == 200
            assert resp2.json()["cached"] is True
            assert fake_client.messages.create.call_count == first_calls

            # Rewrite the file so its mtime moves, then expect a fresh run.
            _make_pptx(pptx_path, slide_count=1)
            # Ensure mtime actually changed on filesystems with low resolution.
            import os
            stat = pptx_path.stat()
            os.utime(pptx_path, ns=(stat.st_atime_ns, stat.st_mtime_ns + 1_000_000))

            resp3 = await client.post(
                "/api/files/beautify-deck", json={"path": "deck.pptx"}
            )
            assert resp3.status_code == 200
            assert resp3.json()["cached"] is False
            assert fake_client.messages.create.call_count > first_calls


@pytest.mark.asyncio
async def test_beautify_requires_api_key(client):
    with tempfile.TemporaryDirectory() as tmpdir:
        tmppath = Path(tmpdir)
        pptx_path = tmppath / "deck.pptx"
        _make_pptx(pptx_path, slide_count=1)

        with patch("routers.projects._projects_root", return_value=tmppath), patch(
            "services.ai_backend.get_ai_client", AsyncMock(return_value=None)
        ), patch.dict("routers.beautify._cache", {}, clear=True):
            resp = await client.post(
                "/api/files/beautify-deck", json={"path": "deck.pptx"}
            )
    assert resp.status_code == 400
    assert "api key" in resp.json()["detail"].lower()
