"""File preview endpoints for the Files page.

Provides rich previews for documents, spreadsheets, slides, and PDFs by
converting them server-side into a renderable representation (HTML, JSON,
or extracted text). Plain text and images are handled by the existing
`/files/read` endpoint in the projects router.

All paths are resolved through the same readable-path helper used by the
projects router so previews stay scoped to the two roots that Recent
Documents walks (the workspace and ``~/.myos/files/``) and cannot
escape them.
"""

from __future__ import annotations

import asyncio
import base64
import collections
import csv
import io
from pathlib import Path

from fastapi import APIRouter, HTTPException, Query
from fastapi.responses import Response

from routers.projects import _resolve_readable_path
from services import recent_deletes

router = APIRouter(tags=["files"])


# Maximum number of spreadsheet rows to return in a preview.
MAX_PREVIEW_ROWS = 500

# Maximum number of columns to return per sheet.
MAX_PREVIEW_COLS = 50

# Maximum PDF pages to extract text from in a single preview.
MAX_PDF_PAGES = 50

# Module-level mtime-keyed LRU cache for parsed PDFs: {(path_str, mtime): result}
_PDF_CACHE: collections.OrderedDict = collections.OrderedDict()
_PDF_CACHE_MAX = 20

# Maximum pptx slides to extract text from.
MAX_PPTX_SLIDES = 200


def _read_docx(path: Path) -> dict:
    """Convert a .docx file to HTML using mammoth."""
    try:
        import mammoth  # type: ignore
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail="Word document preview is not available on this server.",
        ) from exc

    try:
        with path.open("rb") as f:
            result = mammoth.convert_to_html(f)
        html = result.value or ""
    except Exception as exc:
        raise HTTPException(
            status_code=500,
            detail=f"Could not read this Word document: {exc}",
        )

    return {"kind": "docx", "html": html}


def _read_xlsx(path: Path) -> dict:
    """Parse an xlsx file with openpyxl and return sheets as JSON rows."""
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail="Spreadsheet preview is not available on this server.",
        ) from exc

    try:
        wb = load_workbook(path, read_only=True, data_only=True)
    except Exception as exc:
        raise HTTPException(
            status_code=500,
            detail=f"Could not read this spreadsheet: {exc}",
        )

    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[list[object]] = []
        total_rows = 0
        for row in ws.iter_rows(values_only=True):
            total_rows += 1
            if len(rows) >= MAX_PREVIEW_ROWS:
                continue
            trimmed = list(row)[:MAX_PREVIEW_COLS]
            rows.append([_jsonable(cell) for cell in trimmed])

        sheets.append(
            {
                "name": sheet_name,
                "rows": rows,
                "total_rows": total_rows,
                "truncated": total_rows > len(rows),
            }
        )

    wb.close()
    return {"kind": "spreadsheet", "sheets": sheets}


def _read_csv(path: Path) -> dict:
    """Parse a CSV file and return a single sheet of rows."""
    try:
        text = path.read_text(errors="replace")
    except OSError as exc:
        raise HTTPException(status_code=500, detail=f"Could not read file: {exc}")

    reader = csv.reader(io.StringIO(text))
    rows: list[list[str]] = []
    total_rows = 0
    for row in reader:
        total_rows += 1
        if len(rows) >= MAX_PREVIEW_ROWS:
            continue
        rows.append(row[:MAX_PREVIEW_COLS])

    sheet = {
        "name": path.stem or "CSV",
        "rows": rows,
        "total_rows": total_rows,
        "truncated": total_rows > len(rows),
    }
    return {"kind": "spreadsheet", "sheets": [sheet]}


def _read_pptx(path: Path) -> dict:
    """Extract text from each slide of a .pptx file."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail="PowerPoint preview is not available on this server.",
        ) from exc

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise HTTPException(
            status_code=500,
            detail=f"Could not read this presentation: {exc}",
        )

    slides = []
    for idx, slide in enumerate(prs.slides):
        if idx >= MAX_PPTX_SLIDES:
            break
        title = ""
        body_lines: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            try:
                tf = shape.text_frame
                text = (tf.text or "").strip()
            except Exception:
                continue
            if not text:
                continue
            # python-pptx >= 1.0 raises ValueError (not AttributeError) on
            # non-placeholder shapes, so we cannot rely on getattr default.
            try:
                is_placeholder = shape.placeholder_format is not None
            except (ValueError, AttributeError):
                is_placeholder = False
            if not title and is_placeholder:
                title = text.splitlines()[0]
                remainder = "\n".join(text.splitlines()[1:]).strip()
                if remainder:
                    body_lines.append(remainder)
            else:
                body_lines.append(text)
        slides.append(
            {
                "index": idx + 1,
                "title": title,
                "body": "\n".join(body_lines).strip(),
            }
        )

    return {"kind": "slides", "slides": slides}


def _read_pdf(path: Path) -> dict:
    """Extract text from each page of a PDF. Results are mtime-keyed LRU cached."""
    mtime = path.stat().st_mtime
    cache_key = (str(path), mtime)
    if cache_key in _PDF_CACHE:
        _PDF_CACHE.move_to_end(cache_key)
        return _PDF_CACHE[cache_key]

    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail="PDF preview is not available on this server.",
        ) from exc

    try:
        reader = PdfReader(str(path))
    except Exception as exc:
        raise HTTPException(
            status_code=500,
            detail=f"Could not read this PDF: {exc}",
        )

    pages = []
    total_pages = len(reader.pages)
    for idx, page in enumerate(reader.pages):
        if idx >= MAX_PDF_PAGES:
            break
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        pages.append({"index": idx + 1, "text": text.strip()})

    result = {
        "kind": "pdf",
        "pages": pages,
        "total_pages": total_pages,
        "truncated": total_pages > len(pages),
    }
    _PDF_CACHE[cache_key] = result
    _PDF_CACHE.move_to_end(cache_key)
    while len(_PDF_CACHE) > _PDF_CACHE_MAX:
        _PDF_CACHE.popitem(last=False)
    return result


def _jsonable(value: object) -> object:
    """Coerce spreadsheet cell values into JSON-safe forms."""
    if value is None:
        return None
    if isinstance(value, (str, int, float, bool)):
        return value
    # datetimes, Decimal, etc.
    return str(value)


# Extension to handler mapping for rich previews.
_RICH_HANDLERS = {
    ".docx": _read_docx,
    ".xlsx": _read_xlsx,
    ".xlsm": _read_xlsx,
    ".csv": _read_csv,
    ".pptx": _read_pptx,
    ".pdf": _read_pdf,
}


@router.get("/files/preview")
async def preview_file(path: str = Query(..., description="Relative path to the file")):
    """Return a rich preview representation for supported file types.

    For .docx the response includes rendered HTML. For spreadsheets it includes
    one or more sheets of parsed rows. For .pptx it returns slide text. For
    PDFs it returns extracted text per page plus a hint to use the raw bytes
    endpoint for client-side rendering.

    Text and image files are better served by `/files/read` (existing endpoint)
    and `/files/raw` respectively.
    """
    resolved = _resolve_readable_path(path)
    if not resolved.is_file():
        raise HTTPException(status_code=400, detail="Path is not a file.")

    ext = resolved.suffix.lower()
    handler = _RICH_HANDLERS.get(ext)
    if handler is None:
        raise HTTPException(
            status_code=400,
            detail=f"No rich preview available for {ext or 'this file type'}.",
        )

    if ext == ".pdf":
        data = await asyncio.to_thread(handler, resolved)
    else:
        data = handler(resolved)
    data["name"] = resolved.name
    data["size"] = resolved.stat().st_size
    return data


# MIME types for raw byte responses (images, PDFs).
_RAW_MIME = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".svg": "image/svg+xml",
    ".ico": "image/x-icon",
    ".pdf": "application/pdf",
}


@router.get("/files/raw")
async def raw_file(path: str = Query(..., description="Relative path to the file")):
    """Stream raw file bytes with a guessed Content-Type.

    Used by the frontend for images and for client-side PDF rendering.
    """
    resolved = _resolve_readable_path(path)
    if not resolved.is_file():
        raise HTTPException(status_code=400, detail="Path is not a file.")

    ext = resolved.suffix.lower()
    mime = _RAW_MIME.get(ext, "application/octet-stream")

    try:
        data = resolved.read_bytes()
    except OSError as exc:
        raise HTTPException(status_code=500, detail=f"Could not read file: {exc}")

    return Response(content=data, media_type=mime)


# Small helper exposed for tests: encode bytes to a data URL.
def _to_data_url(mime: str, data: bytes) -> str:
    return f"data:{mime};base64,{base64.b64encode(data).decode('ascii')}"


@router.get("/files/provenance")
async def file_provenance(path: str = Query(..., description="Relative or absolute path to the file")):
    """Return the provenance sidecar for a file, or {} if no sidecar exists.

    Same path resolution as /files/preview — accepts workspace-relative
    paths and absolute paths under ~/.myos/files/.
    """
    from services.provenance import read_sidecar

    try:
        resolved = _resolve_readable_path(path)
    except HTTPException:
        return {}
    result = read_sidecar(resolved)
    return result or {}


def _get_files_dir() -> Path:
    return Path.home() / ".myos" / "files"


@router.get("/files/timeline")
async def files_timeline(limit: int = Query(50, ge=1, le=200)):
    """Reverse-chronological feed of agent outputs and uploads with provenance.

    Scans ``~/.myos/files/`` for user-visible files, reads provenance sidecars,
    and returns them sorted by modified_at descending. Sidecar files
    (.provenance.json) are excluded from the listing.
    """
    import mimetypes
    from datetime import datetime, timezone

    from services.provenance import read_sidecar

    files_dir = _get_files_dir()
    if not files_dir.is_dir():
        return {"files": []}

    entries = []
    for f in files_dir.iterdir():
        if not f.is_file():
            continue
        if f.name.endswith(".provenance.json"):
            continue
        try:
            stat = f.stat()
        except OSError:
            continue

        mime, _ = mimetypes.guess_type(f.name)
        provenance = read_sidecar(f)
        if provenance and "source" not in provenance:
            provenance["source"] = "agent"

        entries.append(
            {
                "id": f.name,
                "name": f.name,
                "path": str(f),
                "size": stat.st_size,
                "modified_at": datetime.fromtimestamp(
                    stat.st_mtime, tz=timezone.utc
                ).isoformat(),
                "mime_type": mime or "application/octet-stream",
                "provenance": provenance,
            }
        )

    entries.sort(key=lambda e: e["modified_at"], reverse=True)
    return {"files": entries[:limit]}


@router.post("/files/open")
async def open_file_native(path: str = Query(..., description="Relative or absolute path to the file")):
    """Open a file in the system's default native app (macOS: open, Linux: xdg-open).

    Uses the same safe-path resolution as other file endpoints.
    """
    import subprocess
    import sys

    resolved = _resolve_readable_path(path)
    if not resolved.is_file():
        raise HTTPException(status_code=400, detail="Path is not a file.")

    cmd = ["open", str(resolved)] if sys.platform == "darwin" else ["xdg-open", str(resolved)]
    try:
        subprocess.Popen(cmd)
    except OSError as exc:
        raise HTTPException(status_code=500, detail=f"Could not open file: {exc}")

    return {"ok": True, "path": str(resolved)}


@router.delete("/files/delete")
async def delete_file(path: str = Query(..., description="Relative path to the file to delete")):
    """Delete a file from the workspace.

    Uses the same safe-path resolution as other file endpoints, so path
    traversal outside the project root is blocked.
    """
    resolved = _resolve_readable_path(path)

    if not resolved.is_file():
        raise HTTPException(status_code=400, detail="Path is not a file.")

    try:
        resolved.unlink()
    except OSError as exc:
        raise HTTPException(
            status_code=500,
            detail=f"Could not delete file: {exc}",
        )

    # Tombstone so a racing re-upload or preview does not resurrect the row.
    recent_deletes.record_id(f"file:{path}")
    return {"ok": True, "deleted": str(resolved.name)}


# ---------------------------------------------------------------------------
# Drive import
# ---------------------------------------------------------------------------

from pydantic import BaseModel


class DriveImportBody(BaseModel):
    drive_file_id: str


@router.post("/files/import-from-drive")
async def import_file_from_drive(body: DriveImportBody):
    """Download a Google Drive file into ~/.myos/files/imports/.

    Writes a sidecar JSON next to the file recording provenance
    (source, drive_file_id, imported_at). No inline editing: this is a
    reference-library import only.
    """
    import asyncio
    import json
    import re
    import time

    from services.google_auth import get_credentials, is_authenticated

    if not is_authenticated():
        raise HTTPException(status_code=401, detail="Not connected to Google Drive.")

    file_id = body.drive_file_id.strip()
    if not file_id:
        raise HTTPException(status_code=400, detail="drive_file_id is required.")

    # Fetch file metadata to get the name.
    try:
        from googleapiclient.discovery import build
        from google.oauth2.credentials import Credentials
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail="Google API client is not available on this server.",
        ) from exc

    def _get_meta():
        tokens = get_credentials()
        creds = Credentials(
            token=tokens.get("access_token"),
            refresh_token=tokens.get("refresh_token"),
            token_uri="https://oauth2.googleapis.com/token",
            client_id=None,
            client_secret=None,
        )
        service = build("drive", "v3", credentials=creds)
        return (
            service.files()
            .get(fileId=file_id, fields="id,name,mimeType,size")
            .execute()
        )

    def _download(meta: dict) -> bytes:
        import io
        from googleapiclient.http import MediaIoBaseDownload

        tokens = get_credentials()
        creds = Credentials(
            token=tokens.get("access_token"),
            refresh_token=tokens.get("refresh_token"),
            token_uri="https://oauth2.googleapis.com/token",
            client_id=None,
            client_secret=None,
        )
        service = build("drive", "v3", credentials=creds)
        mime = meta.get("mimeType", "")

        # Google-native files: export as PDF.
        exportable = {
            "application/vnd.google-apps.document",
            "application/vnd.google-apps.presentation",
            "application/vnd.google-apps.spreadsheet",
            "application/vnd.google-apps.drawing",
        }
        if mime in exportable:
            req = service.files().export_media(fileId=file_id, mimeType="application/pdf")
        else:
            req = service.files().get_media(fileId=file_id)

        buf = io.BytesIO()
        dl = MediaIoBaseDownload(buf, req)
        done = False
        while not done:
            _, done = dl.next_chunk()
        return buf.getvalue()

    try:
        meta = await asyncio.get_event_loop().run_in_executor(None, _get_meta)
    except Exception as exc:
        raise HTTPException(
            status_code=500,
            detail=f"Could not get file info from Drive: {exc}",
        ) from exc

    original_name = meta.get("name", "import")
    mime = meta.get("mimeType", "")

    # Sanitise the filename for local storage.
    safe_stem = re.sub(r"[^\w\-. ]", "_", original_name).strip()[:200] or "import"
    # Google-native exports land as PDF.
    native_types = {
        "application/vnd.google-apps.document",
        "application/vnd.google-apps.presentation",
        "application/vnd.google-apps.spreadsheet",
        "application/vnd.google-apps.drawing",
    }
    if mime in native_types and not safe_stem.lower().endswith(".pdf"):
        safe_stem = safe_stem + ".pdf"

    imports_dir = Path.home() / ".myos" / "files" / "imports"
    imports_dir.mkdir(parents=True, exist_ok=True)

    dest = imports_dir / safe_stem
    # Deduplicate: append counter if the name already exists.
    if dest.exists():
        stem = dest.stem
        suffix = dest.suffix
        counter = 1
        while dest.exists():
            dest = imports_dir / f"{stem} ({counter}){suffix}"
            counter += 1

    try:
        content = await asyncio.get_event_loop().run_in_executor(None, lambda: _download(meta))
    except Exception as exc:
        raise HTTPException(
            status_code=500,
            detail=f"Could not download file from Drive: {exc}",
        ) from exc

    try:
        dest.write_bytes(content)
    except OSError as exc:
        raise HTTPException(
            status_code=500,
            detail=f"Could not write file to disk: {exc}",
        ) from exc

    # Write sidecar provenance JSON.
    sidecar = dest.with_suffix(dest.suffix + ".provenance.json")
    provenance = {
        "source": "google-drive",
        "drive_file_id": file_id,
        "drive_name": original_name,
        "imported_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
    }
    try:
        sidecar.write_text(json.dumps(provenance, indent=2))
    except OSError:
        pass  # Sidecar write is best-effort.

    return {
        "ok": True,
        "path": str(dest),
        "name": dest.name,
        "provenance": provenance,
    }
