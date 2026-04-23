"""Documents CRUD router.

Unified create/edit/save surface for all FCP document types: spreadsheets,
slide decks, drawio diagrams, and PDFs. Documents live in ~/.myos/documents/
with metadata tracked in .index.json and version snapshots in .versions/.

Endpoints:
- GET    /api/documents              -- list all documents
- POST   /api/documents              -- create a new document
- POST   /api/documents/upload       -- upload a file from disk
- GET    /api/documents/{id}         -- get document metadata
- DELETE /api/documents/{id}         -- delete a document
- POST   /api/documents/{id}/edit    -- edit via natural-language instruction
- POST   /api/documents/{id}/save    -- save and create version snapshot
- POST   /api/documents/{id}/open-native  -- open in native app
- GET    /api/documents/{id}/versions     -- list version snapshots
- POST   /api/documents/{id}/restore/{version}  -- restore a version
- GET    /api/documents/{id}/preview      -- renderable content for browser display
- GET    /api/documents/{id}/download     -- file download with correct MIME type
"""

from __future__ import annotations

import base64
import json
import logging
import os
import shutil
import subprocess
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional

from fastapi import APIRouter, HTTPException, UploadFile, File
from fastapi.responses import FileResponse
from pydantic import BaseModel, Field

from services.fcp_session import (
    FCPSession,
    FCPSessionError,
    FCPSessionRegistry,
    session_registry,
    DOCUMENTS_DIR,
    VERSIONS_DIR,
)
from services.nl_translator import translate

logger = logging.getLogger(__name__)

router = APIRouter(tags=["documents"])

# ---------------------------------------------------------------------------
# Storage paths
# ---------------------------------------------------------------------------

INDEX_PATH = DOCUMENTS_DIR / ".index.json"

# Map doc type to file extension
_TYPE_TO_EXT: dict[str, str] = {
    "sheet": "xlsx",
    "slides": "pptx",
    "drawio": "drawio",
    "pdf": "pdf",
}

# Map file extension to doc type (reverse of _TYPE_TO_EXT, used for uploads)
_EXT_TO_TYPE: dict[str, str] = {v: k for k, v in _TYPE_TO_EXT.items()}

# Map doc type to FCP server name
_TYPE_TO_SERVER: dict[str, str] = {
    "sheet": "fcp-sheets",
    "slides": "fcp-slides",
    "drawio": "@ostk-ai/fcp-drawio",
    "pdf": "fcp-pdf",
}

# Map doc type to native macOS app name
_TYPE_TO_APP: dict[str, str] = {
    "sheet": "Microsoft Excel",
    "slides": "Microsoft PowerPoint",
    "drawio": "draw.io",
    "pdf": "Preview",
}

# Map extension to MIME type for downloads
_EXT_TO_MIME: dict[str, str] = {
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "drawio": "application/xml",
    "pdf": "application/pdf",
}

# Maximum upload size: 50 MB
_MAX_UPLOAD_BYTES = 50 * 1024 * 1024


# ---------------------------------------------------------------------------
# Request / Response models
# ---------------------------------------------------------------------------

class CreateDocumentRequest(BaseModel):
    type: str = Field(..., description="Document type: sheet, slides, drawio, or pdf")
    name: str = Field(..., min_length=1, max_length=200, description="Document name (without extension)")


class EditDocumentRequest(BaseModel):
    instruction: str = Field(..., min_length=1, description="Natural-language editing instruction")


# ---------------------------------------------------------------------------
# Index helpers
# ---------------------------------------------------------------------------

def _ensure_dirs() -> None:
    """Create the documents directory and versions subdirectory."""
    DOCUMENTS_DIR.mkdir(parents=True, exist_ok=True)
    VERSIONS_DIR.mkdir(parents=True, exist_ok=True)


def _load_index() -> list[dict]:
    """Load the document index from disk. Returns empty list on missing/corrupt."""
    if not INDEX_PATH.exists():
        return []
    try:
        data = json.loads(INDEX_PATH.read_text())
        if isinstance(data, list):
            return data
        return []
    except (json.JSONDecodeError, OSError):
        return []


def _save_index(index: list[dict]) -> None:
    """Persist the document index to disk."""
    _ensure_dirs()
    INDEX_PATH.write_text(json.dumps(index, indent=2))


def _find_doc(doc_id: str, index: list[dict]) -> dict | None:
    """Find a document by ID in the index."""
    for doc in index:
        if doc.get("id") == doc_id:
            return doc
    return None


def _now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


# ---------------------------------------------------------------------------
# Initial content helpers
# ---------------------------------------------------------------------------

def _initial_mutations(doc_type: str, name: str) -> list[dict]:
    """Return the initial FCP mutations for a freshly created document."""
    if doc_type == "sheet":
        return [
            {"verb": "add_sheet", "args": {"name": name}},
            {"verb": "set_cell", "args": {"sheet": name, "cell": "A1", "value": ""}},
        ]
    elif doc_type == "slides":
        return [
            {"verb": "add_slide", "args": {"layout": "title"}},
            {"verb": "set_title", "args": {"text": name}},
        ]
    elif doc_type == "drawio":
        return [
            {"verb": "set_page", "args": {"title": name}},
        ]
    elif doc_type == "pdf":
        return [
            {"verb": "compose", "args": {
                "theme": "modern",
                "blocks": [
                    {"type": "cover", "title": name, "date": _now_iso()[:10]},
                ],
            }},
        ]
    return []


# ---------------------------------------------------------------------------
# Version helpers
# ---------------------------------------------------------------------------

def _version_prefix(doc_id: str) -> str:
    """Return the filename prefix for version snapshots of a document."""
    return f"{doc_id}-v"


def _list_versions(doc_id: str) -> list[dict]:
    """List all version snapshots for a document, sorted by timestamp."""
    _ensure_dirs()
    prefix = _version_prefix(doc_id)
    versions = []
    for f in VERSIONS_DIR.iterdir():
        if f.name.startswith(prefix) and f.is_file():
            # Parse version number from filename: {doc_id}-v{N}-{timestamp}.{ext}
            stem = f.stem  # without final extension
            parts = stem[len(prefix):]  # e.g. "1-20260415T120000"
            version_parts = parts.split("-", 1)
            version_num = version_parts[0] if version_parts else "0"
            ts = version_parts[1] if len(version_parts) > 1 else ""
            versions.append({
                "version": version_num,
                "timestamp": ts,
                "path": str(f),
                "size_bytes": f.stat().st_size,
                "filename": f.name,
            })
    versions.sort(key=lambda v: v.get("version", "0"))
    return versions


def _next_version_number(doc_id: str) -> int:
    """Return the next version number for a document."""
    versions = _list_versions(doc_id)
    if not versions:
        return 1
    try:
        return max(int(v["version"]) for v in versions) + 1
    except (ValueError, KeyError):
        return len(versions) + 1


def _create_version_snapshot(doc: dict) -> dict:
    """Create a version snapshot of the current document file.

    Returns version metadata dict.
    """
    _ensure_dirs()
    file_path = Path(doc["path"])
    if not file_path.exists():
        raise HTTPException(
            status_code=500,
            detail=(
                f"Document file not found at {doc['path']}. "
                f"The file may have been moved or deleted outside of myOS."
            ),
        )

    version_num = _next_version_number(doc["id"])
    ts = datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%S")
    ext = doc.get("ext", _TYPE_TO_EXT.get(doc["type"], "dat"))
    snapshot_name = f"{doc['id']}-v{version_num}-{ts}.{ext}"
    snapshot_path = VERSIONS_DIR / snapshot_name

    shutil.copy2(str(file_path), str(snapshot_path))

    return {
        "version": str(version_num),
        "timestamp": ts,
        "path": str(snapshot_path),
        "size_bytes": snapshot_path.stat().st_size,
        "filename": snapshot_name,
    }


# ---------------------------------------------------------------------------
# Upload helper
# ---------------------------------------------------------------------------

def _resolve_upload_path(original_filename: str) -> Path:
    """Return a collision-free path in DOCUMENTS_DIR for the uploaded file.

    If a file with the same name already exists, a timestamp is appended
    before the extension so the original is never overwritten.
    """
    _ensure_dirs()
    dest = DOCUMENTS_DIR / original_filename
    if not dest.exists():
        return dest

    stem = dest.stem
    suffix = dest.suffix
    ts = datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%S")
    return DOCUMENTS_DIR / f"{stem}-{ts}{suffix}"


# ---------------------------------------------------------------------------
# Endpoints
# ---------------------------------------------------------------------------

@router.get("/documents")
async def list_documents():
    """List all documents from the index.

    Returns a list of document metadata objects.
    """
    index = _load_index()
    return index


@router.post("/documents", status_code=201)
async def create_document(req: CreateDocumentRequest):
    """Create a new document.

    Creates the file on disk, starts an FCP session, writes initial
    content, and adds the document to the index.

    Returns:
        {id, name, type, path}
    """
    doc_type = req.type.lower()
    if doc_type not in _TYPE_TO_EXT:
        raise HTTPException(
            status_code=400,
            detail=(
                f"Unknown document type '{req.type}'. "
                f"Choose one of: {', '.join(sorted(_TYPE_TO_EXT.keys()))}"
            ),
        )

    _ensure_dirs()

    doc_id = str(uuid.uuid4())
    ext = _TYPE_TO_EXT[doc_type]
    # Sanitize filename
    safe_name = req.name.replace("/", "-").replace("\\", "-").strip()
    if not safe_name:
        safe_name = "untitled"
    file_path = DOCUMENTS_DIR / f"{safe_name}.{ext}"

    # Avoid name collisions
    counter = 1
    while file_path.exists():
        file_path = DOCUMENTS_DIR / f"{safe_name}-{counter}.{ext}"
        counter += 1

    now = _now_iso()
    doc = {
        "id": doc_id,
        "name": req.name,
        "type": doc_type,
        "ext": ext,
        "path": str(file_path),
        "created": now,
        "updated": now,
    }

    # Start FCP session and write initial content.
    # Non-fatal: if the FCP server is not installed the document is still
    # created and saved to the index so it appears in the list.
    server_name = _TYPE_TO_SERVER[doc_type]
    try:
        session = session_registry.create(doc_id, server_name, file_path)
        mutations = _initial_mutations(doc_type, req.name)
        try:
            session.mutate(mutations)
            session.save(str(file_path))
        except (FCPSessionError, Exception) as e:
            session_registry.close(doc_id)
            logger.warning(
                "Could not initialize document content for %s (%s): %s",
                doc_id, doc_type, e,
            )
    except (FCPSessionError, Exception) as e:
        logger.warning(
            "Could not start FCP session for %s (%s): %s — document saved without live session",
            doc_id, doc_type, e,
        )

    # Add to index
    index = _load_index()
    index.append(doc)
    _save_index(index)

    return {
        "id": doc_id,
        "name": req.name,
        "type": doc_type,
        "path": str(file_path),
    }


@router.post("/documents/upload", status_code=201)
async def upload_document(file: UploadFile = File(...)):
    """Upload a file and add it to the documents index.

    Accepts multipart form data. Detects the document type from the file
    extension:
      .xlsx -> sheet, .pptx -> slides, .drawio -> drawio, .pdf -> pdf.
    Unknown extensions are stored as a generic file (no FCP session).

    Size limit: 50 MB.

    Returns:
        {id, name, type, path, size_bytes}
    """
    # Read all bytes (streaming would be better for huge files, but the
    # 50 MB cap keeps memory reasonable)
    contents = await file.read()
    size = len(contents)

    if size > _MAX_UPLOAD_BYTES:
        raise HTTPException(
            status_code=413,
            detail="File is too large. Maximum size is 50 MB.",
        )

    original_filename = file.filename or "upload"
    # Sanitize: strip path separators
    original_filename = original_filename.replace("/", "-").replace("\\", "-").strip()
    if not original_filename:
        original_filename = "upload"

    # Detect type from extension
    ext = Path(original_filename).suffix.lstrip(".").lower()
    doc_type = _EXT_TO_TYPE.get(ext, "file")

    _ensure_dirs()

    # Resolve collision-free destination
    dest = _resolve_upload_path(original_filename)

    # Write bytes to disk
    dest.write_bytes(contents)

    doc_id = str(uuid.uuid4())
    now = _now_iso()
    name = Path(original_filename).stem

    doc = {
        "id": doc_id,
        "name": name,
        "type": doc_type,
        "ext": ext if ext else "",
        "path": str(dest),
        "size_bytes": size,
        "source": "upload",
        "created": now,
        "updated": now,
    }

    # For supported FCP types, try to start a session so chat editing
    # is immediately available
    if doc_type in _TYPE_TO_SERVER:
        server_name = _TYPE_TO_SERVER[doc_type]
        try:
            session_registry.create(doc_id, server_name, dest)
        except Exception as e:
            # Non-fatal: the file is saved, just no live session
            logger.warning(
                "Could not start FCP session for uploaded %s (%s): %s",
                original_filename, doc_type, e,
            )

    # Add to index
    index = _load_index()
    index.append(doc)
    _save_index(index)

    return {
        "id": doc_id,
        "name": name,
        "type": doc_type,
        "path": str(dest),
        "size_bytes": size,
    }


@router.get("/documents/{doc_id}")
async def get_document(doc_id: str):
    """Get document metadata and session status.

    Returns the document metadata plus whether an FCP session is active.
    """
    index = _load_index()
    doc = _find_doc(doc_id, index)
    if not doc:
        raise HTTPException(
            status_code=404,
            detail=(
                f"Document '{doc_id}' not found. "
                f"Check the document ID and try again, or list all documents with GET /api/documents."
            ),
        )

    # Check for active session
    session = session_registry.get(doc_id)
    result = dict(doc)
    result["session_active"] = session is not None and session.is_alive

    return result


@router.delete("/documents/{doc_id}", status_code=204)
async def delete_document(doc_id: str):
    """Delete a document.

    Removes the file from disk, closes any active session, and removes
    from the index. Version snapshots are also deleted.
    """
    index = _load_index()
    doc = _find_doc(doc_id, index)
    if not doc:
        raise HTTPException(
            status_code=404,
            detail=(
                f"Document '{doc_id}' not found. "
                f"It may have already been deleted."
            ),
        )

    # Close session if active
    session_registry.close(doc_id)

    # Delete the file
    file_path = Path(doc["path"])
    if file_path.exists():
        try:
            file_path.unlink()
        except OSError as e:
            logger.warning("Could not delete document file %s: %s", file_path, e)

    # Delete version snapshots
    for v in _list_versions(doc_id):
        try:
            Path(v["path"]).unlink()
        except OSError:
            pass

    # Remove from index
    new_index = [d for d in index if d.get("id") != doc_id]
    _save_index(new_index)

    return None


@router.post("/documents/{doc_id}/edit")
async def edit_document(doc_id: str, req: EditDocumentRequest):
    """Edit a document using a natural-language instruction.

    Translates the instruction into FCP mutations using keyword matching
    and applies them to the active session.

    Returns:
        {digest, mutations_applied}
    """
    index = _load_index()
    doc = _find_doc(doc_id, index)
    if not doc:
        raise HTTPException(
            status_code=404,
            detail=(
                f"Document '{doc_id}' not found. "
                f"Check the document ID and try again."
            ),
        )

    # Translate instruction
    result = translate(req.instruction, doc["type"])
    if "error" in result:
        raise HTTPException(
            status_code=400,
            detail=result["error"],
        )

    mutations = result["mutations"]

    # Get or create session
    session = session_registry.get(doc_id)
    if session is None:
        server_name = _TYPE_TO_SERVER[doc["type"]]
        file_path = Path(doc["path"])
        try:
            session = session_registry.create(doc_id, server_name, file_path)
        except (FCPSessionError, Exception) as e:
            raise HTTPException(
                status_code=500,
                detail=(
                    f"Could not start FCP session: {e}. "
                    f"Check that the FCP server is installed and try again."
                ),
            )

    # Apply mutations
    try:
        mutate_result = session.mutate(mutations)
    except (FCPSessionError, Exception) as e:
        raise HTTPException(
            status_code=500,
            detail=(
                f"Failed to apply edit: {e}. "
                f"The instruction was understood but the FCP server rejected it. "
                f"Try a simpler instruction."
            ),
        )

    # Update timestamp in index
    doc["updated"] = _now_iso()
    _save_index(index)

    # Query digest if possible
    digest = None
    try:
        digest = session.query()
    except (FCPSessionError, Exception):
        pass

    return {
        "digest": digest,
        "mutations_applied": len(mutations),
    }


@router.post("/documents/{doc_id}/save")
async def save_document(doc_id: str):
    """Save the document and create a version snapshot.

    Returns:
        {path, size_bytes, version}
    """
    index = _load_index()
    doc = _find_doc(doc_id, index)
    if not doc:
        raise HTTPException(
            status_code=404,
            detail=(
                f"Document '{doc_id}' not found. "
                f"Check the document ID and try again."
            ),
        )

    # Save via session if active
    session = session_registry.get(doc_id)
    if session is not None:
        try:
            session.save(doc["path"])
        except (FCPSessionError, Exception) as e:
            raise HTTPException(
                status_code=500,
                detail=(
                    f"Failed to save: {e}. "
                    f"The FCP session may have expired. Try editing again."
                ),
            )

    # Create version snapshot
    version_info = _create_version_snapshot(doc)

    # Update timestamp
    doc["updated"] = _now_iso()
    _save_index(index)

    file_path = Path(doc["path"])
    size = file_path.stat().st_size if file_path.exists() else 0

    return {
        "path": doc["path"],
        "size_bytes": size,
        "version": version_info["version"],
    }


@router.post("/documents/{doc_id}/open-native")
async def open_native(doc_id: str):
    """Open the document in its native macOS application.

    Returns:
        {opened: true, app: str}
    """
    index = _load_index()
    doc = _find_doc(doc_id, index)
    if not doc:
        raise HTTPException(
            status_code=404,
            detail=(
                f"Document '{doc_id}' not found. "
                f"Check the document ID and try again."
            ),
        )

    file_path = Path(doc["path"])
    if not file_path.exists():
        raise HTTPException(
            status_code=400,
            detail=(
                f"Document file not found at {doc['path']}. "
                f"The file may have been moved or deleted."
            ),
        )

    app_name = _TYPE_TO_APP.get(doc["type"], "Preview")
    try:
        subprocess.Popen(
            ["open", "-a", app_name, str(file_path)],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
    except OSError as e:
        raise HTTPException(
            status_code=500,
            detail=(
                f"Could not open '{app_name}': {e}. "
                f"Make sure {app_name} is installed, or open the file manually at {doc['path']}."
            ),
        )

    return {"opened": True, "app": app_name}


@router.get("/documents/{doc_id}/versions")
async def list_versions(doc_id: str):
    """List all version snapshots for a document.

    Returns a list of version metadata objects.
    """
    index = _load_index()
    doc = _find_doc(doc_id, index)
    if not doc:
        raise HTTPException(
            status_code=404,
            detail=(
                f"Document '{doc_id}' not found. "
                f"Check the document ID and try again."
            ),
        )

    versions = _list_versions(doc_id)
    return versions


@router.post("/documents/{doc_id}/restore/{version}")
async def restore_version(doc_id: str, version: str):
    """Restore a document to a previous version.

    Copies the version snapshot back over the current file and reopens
    the FCP session.

    Returns:
        {restored: true, version: str, path: str}
    """
    index = _load_index()
    doc = _find_doc(doc_id, index)
    if not doc:
        raise HTTPException(
            status_code=404,
            detail=(
                f"Document '{doc_id}' not found. "
                f"Check the document ID and try again."
            ),
        )

    # Find the requested version
    versions = _list_versions(doc_id)
    target = None
    for v in versions:
        if v["version"] == version:
            target = v
            break

    if target is None:
        available = [v["version"] for v in versions]
        raise HTTPException(
            status_code=404,
            detail=(
                f"Version '{version}' not found for document '{doc_id}'. "
                f"Available versions: {available or 'none'}. "
                f"List versions with GET /api/documents/{doc_id}/versions."
            ),
        )

    snapshot_path = Path(target["path"])
    if not snapshot_path.exists():
        raise HTTPException(
            status_code=500,
            detail=(
                f"Version snapshot file not found at {target['path']}. "
                f"The snapshot may have been deleted."
            ),
        )

    # Close existing session
    session_registry.close(doc_id)

    # Copy snapshot over current file
    file_path = Path(doc["path"])
    try:
        shutil.copy2(str(snapshot_path), str(file_path))
    except OSError as e:
        raise HTTPException(
            status_code=500,
            detail=(
                f"Failed to restore version: {e}. "
                f"Check file permissions for {doc['path']}."
            ),
        )

    # Reopen session
    server_name = _TYPE_TO_SERVER[doc["type"]]
    try:
        session_registry.create(doc_id, server_name, file_path)
    except (FCPSessionError, Exception) as e:
        logger.warning("Could not reopen session after restore: %s", e)
        # Restore succeeded even if session reopen fails

    # Update timestamp
    doc["updated"] = _now_iso()
    _save_index(index)

    return {
        "restored": True,
        "version": version,
        "path": doc["path"],
    }


# ---------------------------------------------------------------------------
# Preview helpers
# ---------------------------------------------------------------------------

def _preview_sheet(file_path: Path) -> dict:
    """Extract sheet data from an .xlsx file for browser preview."""
    try:
        import openpyxl
    except ImportError:
        return {
            "type": "sheet",
            "fallback": True,
            "message": (
                "Install openpyxl to see spreadsheet previews: "
                "pip install openpyxl"
            ),
            "size_bytes": file_path.stat().st_size,
        }

    wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
    sheets = []
    for i, ws_name in enumerate(wb.sheetnames[:20]):
        ws = wb[ws_name]
        rows_out: list[list] = []
        headers: list[str] = []
        for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
            str_row = [str(c) if c is not None else "" for c in row]
            if row_idx == 0:
                headers = str_row
            else:
                rows_out.append(str_row)
            if row_idx >= 100:  # cap at 100 data rows + header
                break
        sheets.append({"name": ws_name, "headers": headers, "rows": rows_out})
    wb.close()
    return {"type": "sheet", "sheets": sheets}


def _preview_pdf(file_path: Path) -> dict:
    """Extract text (and optional thumbnails) from a PDF file."""
    try:
        import fitz  # pymupdf
    except ImportError:
        return {
            "type": "pdf",
            "fallback": True,
            "message": (
                "Install PyMuPDF to see PDF previews: "
                "pip install pymupdf"
            ),
            "size_bytes": file_path.stat().st_size,
        }

    doc = fitz.open(str(file_path))
    page_count = len(doc)
    pages = []
    for i, page in enumerate(doc):
        if i >= 20:
            break
        text = page.get_text()
        entry: dict = {"page": i + 1, "text": text}
        try:
            pix = page.get_pixmap(dpi=72)
            img_bytes = pix.tobytes("png")
            b64 = base64.b64encode(img_bytes).decode("ascii")
            entry["thumbnail"] = f"data:image/png;base64,{b64}"
        except Exception:
            pass  # thumbnails are optional
        pages.append(entry)
    doc.close()
    return {"type": "pdf", "page_count": page_count, "pages": pages}


def _preview_drawio(file_path: Path) -> dict:
    """Return drawio XML for browser rendering."""
    xml_content = file_path.read_text(encoding="utf-8", errors="replace")
    return {"type": "drawio", "xml": xml_content}


def _preview_slides(file_path: Path) -> dict:
    """Extract slide text from a .pptx file."""
    try:
        from pptx import Presentation
    except ImportError:
        return {
            "type": "slides",
            "fallback": True,
            "message": (
                "Install python-pptx to see slide previews: "
                "pip install python-pptx"
            ),
            "size_bytes": file_path.stat().st_size,
        }

    prs = Presentation(str(file_path))
    slides_out = []
    for i, slide in enumerate(prs.slides):
        if i >= 50:
            break
        title = ""
        body_parts: list[str] = []
        notes = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text_frame.text
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if text and text != title:
                    body_parts.append(text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
        slides_out.append({
            "number": i + 1,
            "title": title,
            "body": "\n".join(body_parts),
            "notes": notes,
        })
    return {"type": "slides", "slides": slides_out}


_PREVIEW_HANDLERS: dict[str, callable] = {
    "sheet": _preview_sheet,
    "pdf": _preview_pdf,
    "drawio": _preview_drawio,
    "slides": _preview_slides,
}


# ---------------------------------------------------------------------------
# Preview and download endpoints
# ---------------------------------------------------------------------------

@router.get("/documents/{doc_id}/preview")
async def preview_document(doc_id: str):
    """Return renderable content for a document.

    Response shape depends on document type. Degrades gracefully if
    a required library is missing.
    """
    index = _load_index()
    doc = _find_doc(doc_id, index)
    if not doc:
        raise HTTPException(
            status_code=404,
            detail=(
                f"Document '{doc_id}' not found. "
                f"Check the document ID and try again, or list all documents with GET /api/documents."
            ),
        )

    file_path = Path(doc["path"])
    if not file_path.exists():
        raise HTTPException(
            status_code=404,
            detail=(
                "Document file not found. Try creating a new one."
            ),
        )

    handler = _PREVIEW_HANDLERS.get(doc["type"])
    if handler is None:
        return {
            "type": doc["type"],
            "fallback": True,
            "message": f"Preview not supported for type '{doc['type']}'.",
            "size_bytes": file_path.stat().st_size,
        }

    try:
        return handler(file_path)
    except Exception as e:
        logger.warning("Preview failed for %s: %s", doc_id, e)
        return {
            "type": doc["type"],
            "fallback": True,
            "message": f"Preview failed: {e}. Try downloading the file instead.",
            "size_bytes": file_path.stat().st_size,
        }


@router.get("/documents/{doc_id}/download")
async def download_document(doc_id: str):
    """Serve the document file for browser download.

    Sets Content-Disposition to attachment and the correct MIME type.
    """
    index = _load_index()
    doc = _find_doc(doc_id, index)
    if not doc:
        raise HTTPException(
            status_code=404,
            detail=(
                f"Document '{doc_id}' not found. "
                f"Check the document ID and try again, or list all documents with GET /api/documents."
            ),
        )

    file_path = Path(doc["path"])
    if not file_path.exists():
        raise HTTPException(
            status_code=404,
            detail=(
                "Document file not found. Try creating a new one."
            ),
        )

    ext = doc.get("ext", _TYPE_TO_EXT.get(doc["type"], "dat"))
    mime = _EXT_TO_MIME.get(ext, "application/octet-stream")
    filename = file_path.name

    return FileResponse(
        path=str(file_path),
        media_type=mime,
        filename=filename,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )
