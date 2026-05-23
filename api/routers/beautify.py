"""Beautify slides endpoint.

Takes a PowerPoint file path, extracts each slide's content with python-pptx,
and asks Claude to redesign each slide into a structured layout the frontend
can render as clean HTML. Results are cached by file path and modification
time so re-running on an unchanged file is free.

This router is kept separate from ``files.py`` so it can evolve without
touching the in-flight preview work.
"""

from __future__ import annotations

import json
import time
from pathlib import Path
from typing import Any

from fastapi import APIRouter, HTTPException
from pydantic import BaseModel

from routers.projects import _resolve_safe_path


router = APIRouter(tags=["beautify"])

# Hard cap on the number of slides we will redesign at once. Anything
# larger would blow the Claude context budget and cost real money.
MAX_BEAUTIFY_SLIDES = 30

# Allowed fonts the AI may pair. Kept small so the frontend can preload
# them cleanly.
ALLOWED_FONTS = {
    "Inter",
    "Source Sans",
    "Playfair",
    "Merriweather",
    "JetBrains Mono",
    "Space Grotesk",
}

ALLOWED_LAYOUTS = {
    "title",
    "title_and_content",
    "two_column",
    "image_left",
    "image_right",
    "quote",
    "section_header",
    "comparison",
}

# In-process cache. Keyed by (absolute_path, mtime_ns). Value is the list
# of beautified slides. This is deliberately simple. A restart clears it.
_cache: dict[tuple[str, int], list[dict[str, Any]]] = {}


class BeautifyRequest(BaseModel):
    path: str


def _extract_slides(path: Path) -> list[dict[str, Any]]:
    """Pull title, body, bullets, images, and notes out of each slide."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail="Slide reading is not available on this server.",
        ) from exc

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise HTTPException(
            status_code=400,
            detail="This file does not look like a slide deck we can read.",
        ) from exc

    slides: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides):
        title = ""
        bullets: list[str] = []
        body_text: list[str] = []
        image_count = 0

        for shape in slide.shapes:
            shape_type = getattr(shape, "shape_type", None)
            if shape_type is not None and str(shape_type).lower().startswith("picture"):
                image_count += 1
                continue

            if not getattr(shape, "has_text_frame", False):
                continue

            tf = shape.text_frame
            text = (tf.text or "").strip()
            if not text:
                continue

            is_title_placeholder = False
            try:
                pf = getattr(shape, "placeholder_format", None)
                if pf is not None and pf.idx == 0:
                    is_title_placeholder = True
            except (AttributeError, ValueError):
                is_title_placeholder = False

            if is_title_placeholder and not title:
                title = text.splitlines()[0]
                continue

            for para in tf.paragraphs:
                line = (para.text or "").strip()
                if not line:
                    continue
                if getattr(para, "level", 0) >= 0 and len(tf.paragraphs) > 1:
                    bullets.append(line)
                else:
                    body_text.append(line)

        notes = ""
        try:
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes = ""

        slides.append(
            {
                "index": idx + 1,
                "title": title,
                "bullets": bullets,
                "body": "\n".join(body_text).strip(),
                "image_count": image_count,
                "notes": notes,
            }
        )

    return slides


def _prompt_for_slide(slide: dict[str, Any]) -> str:
    """Build the user prompt we send to Claude for a single slide."""
    bullets = "\n".join(f"- {b}" for b in slide.get("bullets", []))
    body = slide.get("body") or ""
    notes = slide.get("notes") or ""
    image_count = slide.get("image_count") or 0

    return (
        "You are redesigning one slide from a presentation. Keep the "
        "meaning of the content the same, but make it cleaner and easier "
        "to read. Think Stripe marketing pages, not WordArt. Never use "
        "em-dashes in any of your output. Use periods or commas instead.\n\n"
        f"Slide number: {slide.get('index')}\n"
        f"Current title: {slide.get('title') or '(no title)'}\n"
        f"Current bullets:\n{bullets or '(none)'}\n"
        f"Current body text:\n{body or '(none)'}\n"
        f"Images on slide: {image_count}\n"
        f"Speaker notes:\n{notes or '(none)'}\n\n"
        "Return ONLY a single JSON object with these fields and no prose:\n"
        "{\n"
        '  "layout": one of "title", "title_and_content", "two_column", '
        '"image_left", "image_right", "quote", "section_header", "comparison",\n'
        '  "title": string,\n'
        '  "body": list of strings (bullets or short paragraphs) OR an '
        "object for two_column and comparison layouts with fields "
        '"left" and "right" each holding a list of strings,\n'
        '  "accent_color": a hex color like "#3b82f6" that fits the tone,\n'
        '  "font_pairing": { "heading": one of Inter, Source Sans, Playfair, '
        "Merriweather, JetBrains Mono, Space Grotesk, "
        '"body": same allowed list },\n'
        '  "notes": optional short string of speaker note suggestions,\n'
        '  "rationale": one sentence in plain language explaining what changed and why.\n'
        "}\n"
    )


def _parse_ai_json(raw: str) -> dict[str, Any]:
    """Pull the first JSON object out of the model's text response."""
    text = raw.strip()
    # Strip markdown fences if the model wrapped the JSON.
    if text.startswith("```"):
        lines = text.splitlines()
        if lines and lines[0].startswith("```"):
            lines = lines[1:]
        if lines and lines[-1].startswith("```"):
            lines = lines[:-1]
        text = "\n".join(lines).strip()

    start = text.find("{")
    end = text.rfind("}")
    if start == -1 or end == -1 or end <= start:
        raise ValueError("model did not return JSON")
    return json.loads(text[start : end + 1])


def _validate_beautified(data: dict[str, Any]) -> dict[str, Any]:
    """Coerce the model output into something the frontend can trust."""
    layout = str(data.get("layout") or "title_and_content")
    if layout not in ALLOWED_LAYOUTS:
        layout = "title_and_content"

    title = str(data.get("title") or "").strip()

    body = data.get("body")
    if body is None:
        body = []

    accent = str(data.get("accent_color") or "#3b82f6")
    if not (accent.startswith("#") and len(accent) in (4, 7)):
        accent = "#3b82f6"

    pairing = data.get("font_pairing") or {}
    heading_font = str(pairing.get("heading") or "Inter")
    body_font = str(pairing.get("body") or "Source Sans")
    if heading_font not in ALLOWED_FONTS:
        heading_font = "Inter"
    if body_font not in ALLOWED_FONTS:
        body_font = "Source Sans"

    notes = str(data.get("notes") or "")
    rationale = str(data.get("rationale") or "").strip()
    if not rationale:
        rationale = "Cleaner layout and clearer wording."

    return {
        "layout": layout,
        "title": title,
        "body": body,
        "accent_color": accent,
        "font_pairing": {"heading": heading_font, "body": body_font},
        "notes": notes,
        "rationale": rationale,
    }


async def _beautify_one_slide(
    client: Any, slide: dict[str, Any]
) -> dict[str, Any]:
    """Call Claude once for a single slide and parse the response."""
    prompt = _prompt_for_slide(slide)
    try:
        response = await client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=1024,
            messages=[{"role": "user", "content": prompt}],
        )
    except Exception as exc:
        # Fall back to a minimal cleaned-up version so the whole deck
        # does not fail if one slide trips the model.
        return _validate_beautified(
            {
                "layout": "title_and_content",
                "title": slide.get("title") or "Slide",
                "body": slide.get("bullets") or [slide.get("body") or ""],
                "accent_color": "#3b82f6",
                "font_pairing": {"heading": "Inter", "body": "Source Sans"},
                "rationale": f"Could not reach the design model: {exc}",
            }
        )

    text_parts: list[str] = []
    for block in getattr(response, "content", []):
        if getattr(block, "type", "") == "text":
            text_parts.append(getattr(block, "text", ""))
    raw = "".join(text_parts)

    try:
        parsed = _parse_ai_json(raw)
    except (ValueError, json.JSONDecodeError):
        parsed = {
            "layout": "title_and_content",
            "title": slide.get("title") or "Slide",
            "body": slide.get("bullets") or [slide.get("body") or ""],
            "rationale": "The design model returned an unexpected shape, so we kept the original content.",
        }

    return _validate_beautified(parsed)


@router.post("/files/beautify-deck")
async def beautify_deck(request: BeautifyRequest) -> dict[str, Any]:
    """Redesign every slide in a PowerPoint file and return before and after."""
    resolved = _resolve_safe_path(request.path)

    if not resolved.is_file():
        raise HTTPException(status_code=400, detail="That path is not a file we can open.")

    if resolved.suffix.lower() != ".pptx":
        raise HTTPException(
            status_code=400,
            detail="This only works on PowerPoint files that end in .pptx.",
        )

    originals = _extract_slides(resolved)
    if len(originals) == 0:
        raise HTTPException(status_code=400, detail="This file has no slides to polish.")

    if len(originals) > MAX_BEAUTIFY_SLIDES:
        raise HTTPException(
            status_code=400,
            detail=(
                f"This deck has {len(originals)} slides. For now we can only "
                f"polish up to {MAX_BEAUTIFY_SLIDES} at a time. Try a shorter file."
            ),
        )

    mtime_ns = resolved.stat().st_mtime_ns
    cache_key = (str(resolved), mtime_ns)
    if cache_key in _cache:
        return {
            "name": resolved.name,
            "slide_count": len(originals),
            "originals": originals,
            "beautified": _cache[cache_key],
            "cached": True,
        }

    from services.ai_backend import get_ai_client
    client = await get_ai_client()
    if client is None:
        raise HTTPException(
            status_code=400,
            detail=(
                "No AI backend configured. Add an Anthropic API key in Settings, or "
                "sign in to Claude Code with your subscription, to use this feature."
            ),
        )

    beautified: list[dict[str, Any]] = []
    start = time.time()
    for slide in originals:
        result = await _beautify_one_slide(client, slide)
        result["index"] = slide["index"]
        beautified.append(result)

    _cache[cache_key] = beautified

    return {
        "name": resolved.name,
        "slide_count": len(originals),
        "originals": originals,
        "beautified": beautified,
        "cached": False,
        "elapsed_seconds": round(time.time() - start, 2),
    }
